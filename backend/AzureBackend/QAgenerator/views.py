from django.shortcuts import render
from django.http import JsonResponse
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
from QAgenerator.forms import QuizForm, AssignmentForm, NoCodeQuizForm
import requests
from decouple import config
from django.shortcuts import redirect
import PyPDF2
from pptx import Presentation
import json 
import spacy 
import nltk
from nltk.tokenize import word_tokenize
from django.views.decorators.csrf import ensure_csrf_cookie

# Environment variables
SECRET_KEY1 = config("SECRET_KEY1")
SECRET_KEY2 = config("SECRET_KEY2")
RESOURCE_NAME = config("RESOURCE_NAME")
MODEL_NAME = config("MODEL_NAME")

def limit_tokens_in_string(text, max_tokens):
    tokens = word_tokenize(text)
    
    # If the number of tokens is below the maximum limit, return the original text
    if len(tokens) <= max_tokens:
        return text
    
    # Otherwise, return the truncated text
    truncated_tokens = tokens[:max_tokens]
    return ' '.join(truncated_tokens)

@ensure_csrf_cookie
def set_csrf_token(request):
    return HttpResponse(status=204)

def home(request):
    return render(request, 'home.html')

def send_message_to_openai(message_content):
    url = f"https://{RESOURCE_NAME}.openai.azure.com/openai/deployments/{MODEL_NAME}/chat/completions?api-version=2023-05-15"

    headers = {
        "Content-Type": "application/json",
        "api-key": f"{SECRET_KEY1}"
    }

    data = {
        "messages": [
            {"role": "system", "content": "You are a helpful assistant that generates quizzes and assignments for educators."},
            {"role": "user", "content": message_content}
        ]
    }

    response = requests.post(url, headers=headers, json=data)
    return response.json()

def get_quiz_answer_key(response):
    quiz =  response['choices'][0]['message']['content']
    url = f"https://{RESOURCE_NAME}.openai.azure.com/openai/deployments/{MODEL_NAME}/chat/completions?api-version=2023-05-15"
    headers = {
        "Content-Type": "application/json",
        "api-key": f"{SECRET_KEY1}"
    }
    data = {
        "messages": [
            {"role": "system", "content": "You are a highly skilled TA that grades quizzes and assignments for educators."},
            {"role": "user", "content": "Given the following quiz, provide a correct answer for each question asked. Ensure it is correct, especially for code questions. Do not add any excessive explanation, and do not add any questions that are not present in the quiz. Provide only one sentance explaining what would qualify the student for earning full points in the question.\nThe quiz is as follows:\n" + quiz},
            {"role": "assistant", "content": "Here is the answer key for the quiz, with each answer showing the number of points per question:"}
        ]
    }

    response = requests.post(url, headers=headers, json=data)
    return response.json()

def no_code_quiz_form(request):

    if request.method == "POST":
        form = NoCodeQuizForm(request.POST, request.FILES)
        if form.is_valid():
            data = form.cleaned_data
            
            user_message = ""
            if 'uploaded_material' in request.FILES:
                uploaded_file = data['uploaded_material']
                entities = extract_content_from_file(uploaded_file)
                file_data = "".join(map(str, entities))
                file_data = limit_tokens_in_string(file_data,4500) +"\n"
                user_message = "Attached is a chunked form of a document submitted by a professor:\n" + file_data
                user_message += "\n\n Using this chunked data only consider the matierial that will fit the topic explanation and requests from the professor mentioned below. You will be using this data to generate a quiz."
                
            user_message += ".\n"

            user_message += f"Write a {data['difficulty_level']} difficulty level quiz with exactly {data['num_questions']} questions, do not go over this number. The quiz should be in a {data['question_style']} format that ensures that the student is tested on their understanding of the topic and complexities behind the question. The topics for this quiz are {data['topic_explanation']}, so ensure that each topic is covered in the quiz. The difficulty level of the quiz should be {data['difficulty_level']}."
            if data['difficulty_level'] == 'elementary':
                user_message+= " The questions should be simple and straightforward."
            elif data['difficulty_level'] == 'intermediate':
                user_message+= " The questions should be a bit complex and require a bit more thought than first glance."
            elif data['difficulty_level'] == 'advanced':
                user_message+= " The questions should be complex and require a lot of thought."
            
            if data['question_style'] == 'short_answer' or data['question_style'] == 'multiple_choice':
                user_message+= " The types of short answer style questions present in the quiz are ones that challenge the student to elaborate on their answer so that it is clear they understand the answer. These can include fill in the blank, short response, etc."
            

            user_message+= "\n In regards to formatting, don't include the type of question in the question itself. For example, don't say 'Question 1 (syntax)', just say 'Question 1'. Also, don't include the answer in the question itself. For example, don't say 'Question 1: What is the output of the following code? print(1+1) Answer: 2', just say 'Question 1: What is the output of the following code? print(1+1)'. Also do not include any notes from the TA in the quiz. Do not include the answer key."
            user_message+= "\n Double check all questions to ensure that they are correct."
            

            response = send_message_to_openai(user_message)
            answer_key = get_quiz_answer_key(response)

            # html = turn_to_html(response['choices'][0]['message']['content'], answer_key['choices'][0]['message']['content'])

            #return render(request, 'results.html', {"response": response, "answer_key": answer_key, "original_quiz": response['choices'][0]['message']['content'], "html": html})
            return JsonResponse({
                "response": response,
                "answer_key": answer_key,
                "original_quiz": response['choices'][0]['message']['content'],
            })
    else:
        return JsonResponse({"errors": form.errors}, status=400)
    #     form = NoCodeQuizForm()
    
    return JsonResponse({"message": "Method not allowed"}, status=405)
    #return render(request, 'quiz_form.html', {"form": form})

@ensure_csrf_cookie
def quiz_form(request):
    if request.method == "POST":
        form = QuizForm(request.POST, request.FILES)
        if form.is_valid():
            data = form.cleaned_data

            user_message = ""
            if 'uploaded_material' in request.FILES: 
                uploaded_file = data['uploaded_material']
                #topics_to_pick = data['topic_explanation']
                entities = extract_content_from_file(uploaded_file)
                file_data = "".join(map(str, entities))
                file_data = limit_tokens_in_string(file_data,4500) +"\n"
                user_message = "Attached is a chunked form of a document submitted by a professor:\n" + file_data
                user_message += "\n\n Using this chunked data only consider the material that will fit the topic explanation and requests from the professor mentioned below. You will be using this data to generate a quiz."
                
            user_message += ".\n"

            user_message = f"Write a coding quiz in the programming language {data['programming_language']} with exactly {data['num_questions']} questions, do not go over this number. The quiz should be in a {data['question_style']} format that ensures that the student is tested on their understanding of syntax and logic. The topics for this quiz are {data['topic_explanation']}. So ensure that each topic is covered in the quiz. The type of questions that can be on this quiz are: {data['question_type']}. The difficulty level of the quiz should be {data['difficulty_level']}."
            if data['difficulty_level'] == 'elementary':
                easy_q = """def greet(name):
                                return "Hello, " + name + "!"
                            message = greet("Alice")

                            After the above code is executed, what will the variable message contain? 
                        """
                user_message+= " The questions should be simple and straightforward. An example of an elementary question would be: " + easy_q
            elif data['difficulty_level'] == 'intermediate':
                medium_q = """def process_data(nums):
                                return [x for x in nums if x%2==0 and x%3!=0]

                            data = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12]
                            result = process_data(data)

                            After executing the above code, what will the result list will contain? 
                        """
                user_message+= " The questions should be a bit complex and require a bit more thought than first glance. An example of an intermediate question would be: " + medium_q
            elif data['difficulty_level'] == 'advanced':
                hard_q = """def f(x):
                                if x == 0:
                                    return 0
                                else:
                                    return x + f(x-1)

                            What is the value of f(3)?
                        """
                user_message+= " The questions should be complex and require a lot of thought. An example of an advanced question would be: " + hard_q
            
            if data['programming_language'] == 'other':
                user_message += f" The specified language is {data['other_language']}."
            
            if data['question_style'] == 'short_answer' or (data['question_style'] == 'multiple_choice' and data['question_style'] == 'short_answer'):
                user_message+= " The types of short answer style questions present in the quiz. The first is questions that give the student a very small snippet of code in the language selected and they must respond with what it will output. The second is questions that ask for a very small snippet of code in the language that fulfils a request."
            
            user_message+= f"\nThis quiz will be graded for a total of {data['total_points']} points. Please show the total points at the top of the quiz."
            if data['fixed_points_per_question']:
                user_message+= " Each question will be worth the same amount of points."
            else:
                user_message+= " Each question will be worth a different amount of points. The points for each question should vary based on the difficulty of the question, but the total points of the collective questions should not exceed the number of points of the quiz. "
            user_message+= "Please include the number of points each question is worth in the question itself. For example, 'Question 1 (5 Points)'."

            user_message+= "\n In regards to formatting, don't include the type of question in the question itself. For example, don't say 'Question 1 (syntax)', just say 'Question 1'. Also, don't include the answer in the question itself. For example, don't say 'Question 1: What is the output of the following code? print(1+1) Answer: 2', just say 'Question 1: What is the output of the following code? print(1+1)'. Also do not include any notes from the TA in the quiz."
            user_message+= "\n Double check all questions including the code snippets to ensure that they are correct."
            
            response = send_message_to_openai(user_message)
            answer_key = get_quiz_answer_key(response)
            response["type"] == "quiz"
            return JsonResponse({
                "response": response,
                "answer_key": answer_key,
                "original_quiz": response['choices'][0]['message']['content']
            })
        else:
            return JsonResponse({"errors": form.errors}, status=400)
    
    return JsonResponse({"message": "Method not allowed"}, status=405)

@csrf_exempt
def modify_quiz(request):
    if request.method == "POST":
        #parse from the JSON that client sends
        data = json.loads(request.body)
        original_quiz = data['original_quiz']
        modifications = data['modifications']

        # Concatenate the original request with modifications
        new_request = "Given the following quiz:" + original_quiz + "\nPlease make the following modifications, but keep absolutely everything else the same except for the question numbers(if questions are removed).\nModifications:" + modifications

        new_request+= "\nIn regards to formatting, don't include the type of question in the question itself. For example, don't say 'Question 1 (syntax)', just say 'Question 1'. Also, don't include the answer in the question itself. For example, don't say 'Question 1: What is the output of the following code? print(1+1) Answer: 2', just say 'Question 1: What is the output of the following code? print(1+1)'. Also do not include any notes from the TA in the quiz."

        response = send_message_to_openai(new_request)
        response["type"] = "quiz"
        answer_key = get_quiz_answer_key(response)

        return JsonResponse({
            "answer_key": answer_key['choices'][0]['message']['content'],
            "modified_quiz" : response['choices'][0]['message']['content']
        })
    return JsonResponse({"message": "Method not allowed"}, status=405)

@csrf_exempt
def modify_assignment(request):
    if request.method == "POST":
        data = json.loads(request.body)
        original_assignment = data['original_assignment']
        modifications = data['modifications']

        # Concatenate the original request with modifications
        new_request = "Given the following quiz:" + original_assignment + "\nPlease make the following modifications, but keep absolutely everything else the same except for the question numbers(if questions are removed).\nModifications:" + modifications

        new_request+= "\nIn regards to formatting, don't include the type of question in the question itself. For example, don't say 'Question 1 (syntax)', just say 'Question 1'. Also, don't include the answer in the question itself. For example, don't say 'Question 1: What is the output of the following code? print(1+1) Answer: 2', just say 'Question 1: What is the output of the following code? print(1+1)'. Also do not include any notes from the TA in the quiz."

        response = send_message_to_openai(new_request)
        response['type'] = "assignment"
        return JsonResponse({"modified_assignment": response})
    
    return JsonResponse({"message": "Method not allowed"}, status=405)

@ensure_csrf_cookie
def assignment_form(request):
    if request.method == "POST":
        form = AssignmentForm(request.POST, request.FILES)
        if form.is_valid():
            data = form.cleaned_data

            user_message = ""
            if 'uploaded_material' in request.FILES:
                uploaded_file = data['uploaded_material']
                entities = extract_content_from_file(uploaded_file)
                file_data = "".join(map(str, entities))
                file_data = limit_tokens_in_string(file_data,4500) +"\n"
                user_message = "Attached is a chunked form of a document submitted by a professor:\n" + file_data
                user_message += "\n\n Using this chunked data only consider the matierial that will fit the topic explanation and requests from the professor mentioned below. You will be using this data to generate a assignment."
                
            user_message += ".\n"

            user_message = f"Generate an assignment on {data['topic_explanation']} in {data['programming_language']} with the following constraints: {data['constraints']}."

            if data['programming_language'] == 'other':
                user_message += f" The specified language is {data['other_language']}."

            response = send_message_to_openai(user_message)
            response["type"] = "assignment"
            return JsonResponse({"assignment": response})
        else:
            return JsonResponse({"errors": form.errors}, status=400)
    
    return JsonResponse({"message": "Method not allowed"}, status=405)

def extract_content_from_file(uploaded_file):
    content = ""

    # Check the file extension
    file_extension = uploaded_file.name.split('.')[-1].lower()

    if file_extension == "pdf":
        # Extract content from PDF
        pdf_reader = PyPDF2.PdfReader(uploaded_file)
        for page_num in range(len(pdf_reader.pages)):
            content += pdf_reader.pages[page_num].extract_text() + "\n"

    elif file_extension in ["ppt", "pptx"]:
        # Extract content from PPT
        presentation = Presentation(uploaded_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
                    
    entities = extract_entities_from_azure(content)

    return entities

def extract_entities_from_azure(text):
        
    SECRET_KEY3 = config("SECRET_KEY3")
    endpoint = "https://usfcslangservice.cognitiveservices.azure.com/"
    
    entity_recognition_url = endpoint + "/text/analytics/v3.0/keyPhrases"
    
    headers = {
        "Ocp-Apim-Subscription-Key": SECRET_KEY3,
        "Content-Type": "application/json"
    }
    max_chars_per_document = 5120
    chunks = [text[i:i + max_chars_per_document] for i in range(0, len(text), max_chars_per_document)]
    
    documents = [{"id": str(idx), "language": "en", "text": chunk} for idx, chunk in enumerate(chunks)]
    
    MAX_DOCS_PER_REQUEST = 5
    responses = []
    for i in range(0, len(documents), MAX_DOCS_PER_REQUEST):
        batch = documents[i:i+MAX_DOCS_PER_REQUEST]
        response = requests.post(entity_recognition_url, headers=headers, json={"documents": batch})
        responses.extend(response.json().get('documents', []))

    all_entities = [doc.get('keyPhrases', []) for doc in responses]
    
    return all_entities