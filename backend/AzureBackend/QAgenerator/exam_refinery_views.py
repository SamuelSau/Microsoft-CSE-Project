from django.shortcuts import render, redirect
from .forms import QuizRefineryForm
import requests
import pdfplumber  # for extracting text from PDFs
from io import BytesIO 
import requests
from decouple import config
from django.shortcuts import redirect
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.http import HttpResponseServerError
import logging
import PyPDF2
from pptx import Presentation
import json 
import spacy 
import nltk
from nltk.tokenize import word_tokenize
from typing import List
from django.http import HttpResponse

SECRET_KEY1 = config("SECRET_KEY1")
SECRET_KEY2 = config("SECRET_KEY2")
SECRET_KEY3 = config("SECRET_KEY3")
RESOURCE_NAME = config("RESOURCE_NAME")
MODEL_NAME = config("MODEL_NAME")

nltk.download('punkt')


def limit_tokens_in_string(text: str, max_tokens: int) -> str:
    tokens = word_tokenize(text)
    
    # If the number of tokens is below the maximum limit, return the original text
    if len(tokens) <= max_tokens:
        return text
    
    # Otherwise, return the truncated text
    truncated_tokens = tokens[:max_tokens]
    return ' '.join(truncated_tokens)

@csrf_exempt
def extract_text_from_pdf(file):
    # Extract text from a PDF file using pdfplumber and an in-memory file
    with pdfplumber.open(file) as pdf:
        text = " ".join(page.extract_text() for page in pdf.pages)
    return text

@csrf_exempt
def get_variant_answer_key(response):
    print(response)
    quiz =  response['choices'][0]['message']['content']
    url = f"https://{RESOURCE_NAME}.openai.azure.com/openai/deployments/{MODEL_NAME}/chat/completions?api-version=2023-05-15"
    headers = {
        "Content-Type": "application/json",
        "api-key": f"{SECRET_KEY1}"
    }
    data = {
        "messages": [
            {"role": "system", "content": "You are a highly skilled TA. Analyze the provided quiz. Generate answers for each question in the quiz and its variants. Preserve the original quiz format, especially if it's not multiple-choice. If the quiz is incomplete or unusual, ignore and proceed to create answers for the provided questions. Ensure the answer key matches the question formats without modifications."},
            {"role": "user", "content": f"Given the following quiz, provide a correct answer for each question. Ensure accuracy, especially for code questions. The quiz is:\n{quiz}"},
            {"role": "assistant", "content": "Here is the answer key for the quiz, with each answer numbered corresponding to the original and variant questions. Do not include the questions, only the answers."}
        ]
    }
    variantResponse = requests.post(url, headers=headers, json=data)
    return variantResponse.json()

@csrf_exempt
def create_quiz_variations(file_content: str, num_variations: int):
    print(file_content)
    
    url = f"https://{RESOURCE_NAME}.openai.azure.com/openai/deployments/{MODEL_NAME}/chat/completions?api-version=2023-05-15"

    headers = {
        "Content-Type": "application/json",
        "api-key": f"{SECRET_KEY1}"
    }

    data = {
        "messages": [
            {"role": "system", "content": "You are an assistant tasked with creating quiz variations. Maintain the same level of difficulty and concepts as the original, and ensure the format, especially non-multiple-choice, is preserved. Vary question order in each variation."},
            {"role": "user", "content": f"Create {str(num_variations)} variations of this quiz. Keep the format identical to the original, which may not be multiple-choice. Surround any code with '```'. Include the original quiz at the top. The original quiz is as follows:\n" + file_content},
            {"role": "assistant", "content": f"Here are {str(num_variations)} variations of the quiz, each titled 'Variation Number'. The original quiz is included at the top. All code is within '```'. Each variation matches the original quiz in format, question count, difficulty, and concepts covered. No answer key is provided here, only questions."}
        ]
    }


    response = requests.post(url, headers=headers, json=data)
    response_json = response.json()  # Get the JSON content of the response
    answer_key = get_variant_answer_key(response_json)
    combined_data = {
        "response": response_json,
        "answer_key": answer_key
    }

    # Convert the combined dictionary to a JSON string
    return combined_data
    
# Set up logging
logger = logging.getLogger(__name__)

@csrf_exempt
def exam_refine_form(request):
    if request.method == "POST":
        try:
            form = QuizRefineryForm(request.POST, request.FILES)
            if form.is_valid():
                num_variations = form.cleaned_data['num_variations']
                uploaded_file = request.FILES['upload_file']
                file_extension = uploaded_file.name.split('.')[-1].lower()  # Extract file extension here

                file_in_memory = BytesIO(uploaded_file.read())
                file_in_memory.seek(0)  # Reset file pointer to the start
                extracted_text = extract_content_from_file(file_in_memory, file_extension)
                response = create_quiz_variations(extracted_text, num_variations)
                response['type'] = 'variation'
                return JsonResponse({"response": response})
            else:
                return JsonResponse({"status": "Invalid Form", "message": "Form data is not valid."})
        except Exception as e:
            # Log the error
            logger.error('Error in exam_refine_form: %s', str(e))

            # Optionally, return a more descriptive error response
            return HttpResponseServerError('Internal Server Error')
    else:
        return JsonResponse({"status": "Form not submitted", "message": "Please submit the form using POST method."})

def extract_content_from_file(uploaded_file, file_extension) -> str:
    content = ""

    if file_extension == "pdf":
        with pdfplumber.open(uploaded_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:  # Check if page text is not None or empty
                    content += page_text + "\n"
    
    #needs to fixing 
    elif file_extension in ["ppt", "pptx"]:

        # Extract content from PPT
        uploaded_file.seek(0)

        # Extract content from PPT or PPTX
        presentation = Presentation(uploaded_file)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
                
    # entities = extract_entities_from_azure(content)    
    # entities_str = '\n'.join([' '.join(phrase_list) for phrase_list in entities])
    # limited_str = limit_tokens_in_string(entities_str, 4500)
    return content

def extract_entities_from_azure(text: str) -> List[List]:
        
    SECRET_KEY3 = config("SECRET_KEY3")
    endpoint = "https://usfcslangservice.cognitiveservices.azure.com/"
    
    entity_recognition_url = endpoint + "/text/analytics/v3.0/keyPhrases"
    
    headers = {
        "Ocp-Apim-Subscription-Key": SECRET_KEY3,
        "Content-Type": "application/json"
    }
    max_chars_per_document = 5120
    chunks = [text[i:i + max_chars_per_document] for i in range(0, len(text), max_chars_per_document)]
    
    documents = [{"id": str(idx), "language": "en", "text": chunk} for idx, chunk in enumerate(chunks)]
    

    MAX_DOCS_PER_REQUEST = 5
    responses = []
    for i in range(0, len(documents), MAX_DOCS_PER_REQUEST):
        batch = documents[i:i+MAX_DOCS_PER_REQUEST]
        response = requests.post(entity_recognition_url, headers=headers, json={"documents": batch})
        responses.extend(response.json().get('documents', []))

    
    all_entities = [doc.get('keyPhrases', []) for doc in responses]
    
    return all_entities